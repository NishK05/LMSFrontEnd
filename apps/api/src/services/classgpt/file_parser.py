import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


def parse_folder(folder_path: str) -> dict[str, str]:
    """
    Given a folder path, recursively walks through it and parses all 
    .pdf, .docx, .pptx, and .txt files.

    Returns a dictionary where:
      - keys are relative filenames (e.g., "week1/notes.pdf")
      - values are the extracted, cleaned plain text
    """
    supported_exts = {'.pdf', '.docx', '.pptx', '.txt'}
    result = {}
    for root, _, files in os.walk(folder_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext not in supported_exts:
                continue
            rel_path = os.path.relpath(os.path.join(root, file), folder_path)
            abs_path = os.path.join(root, file)
            print(f"Parsing {rel_path}...")
            try:
                if ext == '.pdf':
                    text = extract_pdf(abs_path)
                elif ext == '.docx':
                    text = extract_docx(abs_path)
                elif ext == '.pptx':
                    text = extract_pptx(abs_path)
                elif ext == '.txt':
                    text = extract_txt(abs_path)
                else:
                    continue
                cleaned = text.strip()
                result[rel_path] = cleaned
                print(f"  Done: {len(cleaned)} chars")
            except Exception as e:
                print(f"  Failed to parse {rel_path}: {e}")
    return result

def extract_pdf(path):
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(path):
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_txt(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read() 